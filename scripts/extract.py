# -*- coding: utf-8 -*-
"""课件文本提取器：把 PDF / PPTX / DOCX / TXT / MD 转成纯文本，供整理笔记用。

为什么单独做一个脚本：不同格式的课件需要不同的库，而且 Windows 默认 GBK
编码会让中文报 UnicodeEncodeError——这里统一用 UTF-8 输出，避免每次踩坑。

用法：
    python extract.py 课件/第07讲.pdf                # 提取文字到 stdout
    python extract.py 课件/第07讲.pptx --out out.txt  # 写文件
    python extract.py 课件/                           # 整个目录逐个文件提取
    python extract.py 课件/第07讲.pdf --images img/   # 把每页/每图渲染成 PNG
    python extract.py 课件/第07讲.pdf --scan          # 只报告每页文字量

⚠️ 重要：纯文字提取只能拿到"机器文字"，**图片里的内容（乐谱、示意图、
扫描页、图表中的文字、画作等）抽不出来**。所以整理时是"两条腿走路"：
  1) 先用默认模式抽文字（快、省）；
  2) 对文字量很少或含重要图示的页/张，用 --images 渲染成 PNG，再让 Claude
     "看图"读取（模型能直接读图），确保不漏关键信息。
提取结束时会打印每页文字量统计，文字量低（如 <40 字）的页就是需要看图的。

依赖（按需安装）：pymupdf(PDF)  python-pptx(PPTX)  python-docx(DOCX)  Pillow(图片)
注意：音频/视频一律不抽取、不试听——听辨题的曲名/答案课件文字里通常已
写明，照文字整理即可。
"""
import io
import os
import sys

# Windows 控制台默认 GBK，中文会崩；强制 UTF-8 输出
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')


def extract_pdf(path):
    import fitz  # pymupdf
    doc = fitz.open(path)
    out = []
    for i in range(doc.page_count):
        out.append(f"=== 第 {i+1} 页 / PAGE {i+1} ===")
        out.append(doc[i].get_text())
    return "\n".join(out)


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"=== 第 {i} 张幻灯片 / SLIDE {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        out.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    out.append(" | ".join(cells))
        # 演讲者备注里常有补充讲解，一并取出
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                out.append(f"[备注] {note}")
    return "\n".join(out)


def extract_docx(path):
    import docx
    d = docx.Document(path)
    out = []
    for p in d.paragraphs:
        if p.text.strip():
            out.append(p.text)
    for t in d.tables:
        for row in t.rows:
            out.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(out)


def extract_text(path):
    for enc in ('utf-8', 'gbk', 'utf-16'):
        try:
            with open(path, encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    with open(path, encoding='utf-8', errors='replace') as f:
        return f.read()


# ---------- 文字量统计 & 图片渲染（应对"内容在图里"的情况）----------

def density(path):
    """返回 [(标签, 文字字数, 图片数)]，用来判断哪些页/张需要"看图"读取。"""
    ext = os.path.splitext(path)[1].lower()
    rows = []
    if ext == '.pdf':
        import fitz
        doc = fitz.open(path)
        for i in range(doc.page_count):
            txt = doc[i].get_text().strip()
            imgs = len(doc[i].get_images(full=True))
            rows.append((f"第{i+1}页", len(txt), imgs))
    elif ext in ('.pptx', '.ppt'):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, 1):
            txt = sum(len(s.text_frame.text) for s in slide.shapes if s.has_text_frame)
            imgs = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
            rows.append((f"第{i}张", txt, imgs))
    return rows


def render_images(path, outdir, dpi=150, low_text=None):
    """把页/图渲染成 PNG，供 Claude 看图读取。
    low_text 不为 None 时，只渲染文字量 < low_text 的页（聚焦图片页，省资源）。"""
    os.makedirs(outdir, exist_ok=True)
    ext = os.path.splitext(path)[1].lower()
    stem = os.path.splitext(os.path.basename(path))[0]
    saved = []
    if ext == '.pdf':
        import fitz
        doc = fitz.open(path)
        dens = {r[0]: r[1] for r in density(path)}
        for i in range(doc.page_count):
            if low_text is not None and dens.get(f"第{i+1}页", 0) >= low_text:
                continue
            pix = doc[i].get_pixmap(dpi=dpi)
            fn = os.path.join(outdir, f"{stem}_p{i+1:03d}.png")
            pix.save(fn); saved.append(fn)
    elif ext in ('.pptx', '.ppt'):
        # python-pptx 无法整张渲染幻灯片，更**抓不到公式对象（OMML）和版式**，
        # 这里只导出每张内嵌的图片 blob。公式页/整页版式请改用 pptx_to_pdf.py
        # 先转 PDF 再 --images（见末尾提示）。坏图会被跳过而不是中断整批。
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(path)
        skipped = 0
        for si, slide in enumerate(prs.slides, 1):
            k = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    k += 1
                    try:  # 坏图（PIL UnidentifiedImageError 等）跳过，别让整批崩
                        img = shape.image
                        fn = os.path.join(outdir, f"{stem}_s{si:03d}_{k}.{img.ext}")
                        with open(fn, 'wb') as f:
                            f.write(img.blob)
                        saved.append(fn)
                    except Exception:
                        skipped += 1
        if skipped:
            print(f"[提示] 跳过 {skipped} 张无法读取的内嵌图。")
        print("[提示] pptx 的 --images 只取内嵌图、抓不到公式对象。公式/版式页请用："
              f'\n   python "{os.path.join(os.path.dirname(__file__), "pptx_to_pdf.py")}" "{path}"'
              "\n   转成 PDF 后再对该 PDF --images 整页渲染。")
    return saved


EXTRACTORS = {
    '.pdf': extract_pdf,
    '.pptx': extract_pptx,
    '.ppt': extract_pptx,   # 老式 .ppt 多半失败，提示用户先另存为 .pptx
    '.docx': extract_docx,
    '.doc': extract_docx,
    '.txt': extract_text,
    '.md': extract_text,
}


def extract_one(path):
    ext = os.path.splitext(path)[1].lower()
    fn = EXTRACTORS.get(ext)
    if fn is None:
        return f"[跳过] 不支持的格式：{path}（支持 pdf/pptx/docx/txt/md）"
    try:
        return fn(path)
    except ImportError as e:
        lib = {'fitz': 'pymupdf', 'pptx': 'python-pptx', 'docx': 'python-docx'}
        miss = str(e).split("'")[-2] if "'" in str(e) else str(e)
        return f"[需要安装依赖] {path}: pip install {lib.get(miss, miss)}"
    except Exception as e:
        if ext in ('.ppt', '.doc'):
            return f"[失败] {path}: 旧版 {ext} 格式难以解析，请在 PowerPoint/Word 里另存为 .pptx/.docx 再试。详情：{e}"
        return f"[失败] {path}: {e}"


def density_report(path, low=40):
    rows = density(path)
    if not rows:
        return ""
    low_pages = [r[0] for r in rows if r[1] < low]
    lines = ["", f"---- 文字量统计（{os.path.basename(path)}）----"]
    for label, chars, imgs in rows:
        flag = "  ⚠️文字少，建议看图" if chars < low else ""
        lines.append(f"  {label}: 文字{chars}字, 图{imgs}张{flag}")
    if low_pages:
        lines.append(f"⚠️ 以下页文字量低，可能内容在图里，建议 --images 渲染后看图：{', '.join(low_pages)}")
    return "\n".join(lines)


def main():
    argv = sys.argv[1:]
    args = [a for a in argv if not a.startswith('--')]
    out_path = argv[argv.index('--out') + 1] if '--out' in argv else None
    img_dir = argv[argv.index('--images') + 1] if '--images' in argv else None
    scan_only = '--scan' in argv
    low_only = '--lowtext-only' in argv  # 配合 --images：只渲染文字少的页
    if not args:
        print(__doc__); return
    target = args[0]
    files = []
    if os.path.isdir(target):
        files = [os.path.join(target, f) for f in sorted(os.listdir(target))
                 if os.path.isfile(os.path.join(target, f))]
    else:
        files = [target]

    if scan_only:
        for p in files:
            print(density_report(p))
        return

    if img_dir:
        for p in files:
            saved = render_images(p, img_dir, low_text=(40 if low_only else None))
            print(f"{os.path.basename(p)}: 渲染 {len(saved)} 张图 -> {img_dir}")
            for s in saved:
                print("   ", s)
        print("\n提示：用 Read 工具逐张查看这些 PNG，把图里的关键信息补进笔记。")
        return

    chunks = []
    for p in files:
        if len(files) > 1:
            chunks.append(f"\n########## 文件：{os.path.basename(p)} ##########")
        chunks.append(extract_one(p))
        rep = density_report(p)
        if rep:
            chunks.append(rep)
    result = "\n".join(chunks)
    if out_path:
        with open(out_path, 'w', encoding='utf-8') as f:
            f.write(result)
        print(f"已写入 {out_path}（{len(result)} 字）")
    else:
        print(result)


if __name__ == '__main__':
    main()
